# =========================================================
# BTEC / LMS RAG AI Assistant
# Single File Version
# =========================================================
# FEATURES:
# - Reads approved local files ONLY
# - Supports PDF / DOCX / TXT
# - Creates embeddings
# - Stores vectors in ChromaDB
# - Answers ANY query from approved sources only
# - Prevents hallucination
# - Streamlit UI
# - Local-first architecture
#
# RUN:
# pip install streamlit chromadb sentence-transformers
# pip install pypdf python-docx ollama
#
# START OLLAMA:
# ollama run llama3
#
# RUN APP:
# streamlit run app.py
# =========================================================

import os
import glob
import shutil
import streamlit as st
import chromadb
import requests
import re
import time

from io import BytesIO
from pypdf import PdfReader
from docx import Document
from sentence_transformers import SentenceTransformer
import ollama
from fpdf import FPDF
from pptx import Presentation


# =========================================================
# PAGE CONFIG
# =========================================================

st.set_page_config(
    page_title="BTEC RAG Assistant",
    layout="wide"
)


# =========================================================
# UI STYLE + BACKGROUND
# =========================================================

st.markdown("""
<style>

/* PAGE BACKGROUND IMAGE */
.stApp {
    background-image:
        linear-gradient(rgba(255,255,255,0.88), rgba(255,255,255,0.88)),
        url("ISO.png");
    background-size: cover;
    background-position: center;
    background-repeat: no-repeat;
    background-attachment: fixed;
}


/* MAIN CONTAINER */
.block-container {
    position: relative;
    z-index: 10;
    max-width: 850px;
    margin: auto;
    padding-top: 4rem;
}


/* TITLE */
.title {
    color: black;
    font-size: 44px;
    font-weight: 800;
    text-align: center;
    margin-bottom: 25px;
}


/* LABELS */
label {
    color: black !important;
    font-weight: 700 !important;
}

[data-testid="stWidgetLabel"] p {
    color: black !important;
    font-weight: 700 !important;
    font-size: 18px !important;
}


/* TEXTAREA */
.stTextArea textarea {
    background-color: white !important;
    color: black !important;
    border: 1px solid #ccc !important;
    border-radius: 16px !important;
    padding: 14px !important;
    font-size: 16px !important;
    caret-color: black !important;
}

.stTextArea textarea::placeholder {
    color: #555 !important;
}


/* BUTTON */
.stButton>button {
    width: 100%;
    border-radius: 14px;
    height: 50px;
    background: #2563eb;
    color: white;
    border: none;
    font-weight: 700;
    font-size: 16px;
}

.stButton>button:hover,
.stButton>button:focus,
.stButton>button:active {
    background: #1d4ed8 !important;
    color: white !important;
    border: none !important;
    outline: none !important;
    box-shadow: none !important;
}


/* WARNING MESSAGE */
div[data-baseweb="notification"] {
    color: black !important;
}

div[data-baseweb="notification"] * {
    color: black !important;
}


/* ANSWER BOX */
.answer-box {
    background: white;
    color: black;
    padding: 20px;
    border-radius: 18px;
    box-shadow: 0 4px 20px rgba(0,0,0,0.15);
}


/* SOURCE BOX */
.source-box {
    background: #f8fafc;
    color: black;
    padding: 14px;
    border-radius: 14px;
    margin-bottom: 10px;
    border-left: 4px solid #3b82f6;
}


/* HEADINGS */
h1, h2, h3, h4 {
    color: black !important;
}


/* SIDEBAR */
[data-testid="stSidebar"] {
    background: #0f172a;
}

[data-testid="stSidebar"] * {
    color: white !important;
}


/* SPINNER */
.stSpinner > div {
    color: black !important;
}


/* DOWNLOAD BUTTON */
.stDownloadButton > button {
    background-color: #000000 !important;
    color: white !important;
    border-radius: 12px !important;
    border: none !important;
    font-weight: 700 !important;
    height: 48px !important;
    padding: 0 20px !important;
}

/* hover */
.stDownloadButton > button:hover {
    background-color: #1d4ed8 !important;
    color: white !important;
}

/* click */
.stDownloadButton > button:active,
.stDownloadButton > button:focus {
    background-color: #1e40af !important;
    color: white !important;
    border: none !important;
    outline: none !important;
    box-shadow: none !important;
}

</style>
""", unsafe_allow_html=True)


# =========================================================
# SIDEBAR
# =========================================================

st.sidebar.title("📚 RAG System Info")

st.sidebar.info("""
✔ Retrieval-Augmented Generation (RAG)

✔ Uses ONLY approved BTEC documents

✔ No internet knowledge allowed
""")

st.sidebar.markdown("### ⚙️ Used Systems")
st.sidebar.write("🧠 RAG System")
st.sidebar.write("🗂️ ChromaDB")
st.sidebar.write("🤖 Ollama")
st.sidebar.write("🔗 PDF URL Support")
st.sidebar.write("📄 SentenceTransformers")

st.sidebar.markdown("### 📁 Supported Files")
st.sidebar.write("📄 PDF")
st.sidebar.write("📝 DOCX")
st.sidebar.write("📃 TXT")

st.sidebar.markdown("### 📚 Approved Sources")
st.sidebar.write("✔ Pearson BTEC Specifications")
st.sidebar.write("✔ Official BTEC Documents")
st.sidebar.write("✔ Assessment Guides")


# =========================================================
# CONFIG
# =========================================================

DATA_FOLDER = "data"
CHROMA_DB_PATH = "chroma_db"
COLLECTION_NAME = "btec_lms"

#import shutil

#if os.path.exists(CHROMA_DB_PATH):
  # shutil.rmtree(CHROMA_DB_PATH)

@st.cache_resource
def load_embedding_model():
    return SentenceTransformer("all-MiniLM-L6-v2", device="cpu")

embedding_model = load_embedding_model()
client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
collection = client.get_or_create_collection(name=COLLECTION_NAME)

PDF_URLS = [
    "https://qualifications.pearson.com/content/dam/pdf/btec-international-level-2/business/2022/specification-and-sample-assessments/business-specification.pdf"
]

APPROVED_SOURCES = [
    "https://qualifications.pearson.com/",
    "https://iso-lms.online/moodle/"
]

SIMILARITY_THRESHOLD = 0.50


# =========================================================
# FUNCTIONS
# =========================================================

def detect_unit(text):
    patterns = [
        r"(Unit\s*\d+)",
        r"(Unit\s*:\s*\d+)",
        r"(UNIT\s*\d+)"
    ]

    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            return match.group(1)

    return "Unknown Unit"

def read_pdf(path):
    pages = []

    try:
        if path.startswith("http"):
            r = requests.get(path, headers={"User-Agent": "Mozilla/5.0"})
            pdf = PdfReader(BytesIO(r.content))
        else:
            pdf = PdfReader(path)

        for i, page in enumerate(pdf.pages):
            text = page.extract_text()

            if text:
                pages.append({
                    "text": text,
                    "page": i + 1,
                    "source": path
                })

    except Exception as e:
        print("PDF ERROR:", e)

    return pages


def chunk(text, size=600, overlap=80):
    chunks = []
    start = 0

    while start < len(text):
        chunks.append(text[start:start + size])
        start += size - overlap

    return chunks


def build():
    if collection.count() > 0:
        return

    docs = []

    files = glob.glob(
        os.path.join(DATA_FOLDER, "**/*.*"),
        recursive=True
    )

    for f in files:
        if f.endswith(".pdf"):
            docs.append({"file": f, "pages": read_pdf(f)})

        elif f.endswith(".docx"):
            text = "\n".join([p.text for p in Document(f).paragraphs])
            docs.append({
                "file": f,
                "pages": [{"text": text, "page": "Unknown"}]
            })

        elif f.endswith(".txt"):
            docs.append({
                "file": f,
                "pages": [{
                    "text": open(f, encoding="utf-8").read(),
                    "page": "Unknown"
                }]
            })

    # approved Pearson PDF only once
    docs.append({
        "file": PDF_URLS[0],
        "pages": read_pdf(PDF_URLS[0])
    })

    doc_id = 0

    for d in docs:
        for p in d["pages"]:
            unit = detect_unit(p["text"])

            for c in chunk(p["text"]):
                emb = embedding_model.encode(c).tolist()

                collection.add(
                    ids=[str(doc_id)],
                    documents=[c],
                    embeddings=[emb],
                    metadatas=[{
                        "source": d["file"],
                        "page": p["page"],
                        "unit": unit
                    }]
                )

                doc_id += 1


def retrieve(q):
    q_emb = embedding_model.encode(q).tolist()

    res = collection.query(
        query_embeddings=[q_emb],
        n_results=2
    )

    docs = res["documents"][0]
    metas = res["metadatas"][0]
    dist = res["distances"][0]

    results = []

    # semantic search first
    for d, m, dis in zip(docs, metas, dist):
        similarity = 1 / (1 + dis)

        print("SIMILARITY:", similarity)

        if similarity >= SIMILARITY_THRESHOLD:
            results.append({
                "text": d,
                "source": m.get("source", "Unknown"),
                "page": m.get("page", "Unknown"),
                "unit": m.get("unit", "Unknown"),
                "score": similarity
            })



    return results

def answer(q, task_type):
    ctx = retrieve(q)

    if not ctx:
        return ("Cannot answer based on the available approved sources.", [])

    context = "\n\n".join([c["text"] for c in ctx[:2]])

    prompts = {
    "Normal Answer": f"""
You are a strict BTEC academic assistant.

Use ONLY the provided context.

Context:
{context}

Question:
{q}

Rules:
- Answer only if explicitly found in context.
- Do NOT infer.
- Do NOT add outside knowledge.
- If information is incomplete or missing, reply exactly:
Cannot answer based on the available approved sources.
""",


        "Quiz Questions": f"""
Use ONLY this approved context.

Context:
{context}

Create 5 quiz questions with answers based ONLY on the context.

Format:
1. Question
Answer: ...

Rules:
- No external knowledge.
""",

        "Assignment Guidance": f"""
Use ONLY this approved context.

Context:
{context}

Create assignment guidance for this request:
{q}

Include:
- Task objective
- Key points student should include
- Suggested structure

Rules:
- Only approved content.
""",

        "Lesson Summary": f"""
Use ONLY this approved context.

Context:
{context}

Generate a lesson summary.

Include:
- Main concepts
- Important definitions
- Key takeaways

Rules:
- Only approved content.
""",

        "Learning Activities": f"""
Use ONLY this approved context.

Context:
{context}

Create 3 learning activities based only on context.

Examples:
- discussion
- worksheet
- research task
""",

        "Short Explanation": f"""
Use ONLY this approved context.

Context:
{context}

Provide a short simple explanation for:
{q}

Rules:
- concise
- clear
- approved content only
"""
    }

    res = ollama.chat(
        model="llama3",
        messages=[{
            "role": "user",
            "content": prompts[task_type]
        }]
    )

    final_answer = res["message"]["content"].strip()
    if "Cannot answer based on the available approved sources" not in final_answer:
        if len(final_answer.split()) > 180:
            final_answer = "Cannot answer based on the available approved sources."
    return final_answer, ctx
# =========================================================
# BUILD DATABASE
# =========================================================

@st.cache_resource
def initialize_database():
    build()

initialize_database()

#print("DATABASE BUILT")
print("\n========= DATABASE REVIEW =========")

all_docs = collection.get(include=["documents", "metadatas"])

print("TOTAL DOCUMENT CHUNKS:", len(all_docs["documents"]))

for i in range(min(3, len(all_docs["documents"]))):
    print(f"\n--- SAMPLE {i+1} ---")
    print("SOURCE:", all_docs["metadatas"][i]["source"])
    print("PAGE:", all_docs["metadatas"][i]["page"])
    print("UNIT:", all_docs["metadatas"][i]["unit"])
    print("TEXT SAMPLE:")
    print(all_docs["documents"][i][:500])

print("\n===================================")

def create_pdf(text):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font("Arial", "B", 16)
    pdf.cell(0, 10, "BTEC Generated Content", ln=True, align="C")
    pdf.ln(10)

    pdf.set_font("Arial", size=12)

    lines = text.split("\n")
    lines = [line.strip() for line in lines if line.strip()]

    for line in lines:
        pdf.multi_cell(0, 10, line)
        pdf.ln(2)

    file_name = "output.pdf"
    pdf.output(file_name)
    return file_name

def create_ppt(text):
    prs = Presentation()

    lines = text.split("\n")

    # إزالة الأسطر الفارغة
    lines = [line.strip() for line in lines if line.strip()]

    chunk_size = 5   # عدد الأسطر بكل سلايد

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BTEC Generated Content"
    slide.placeholders[1].text = "Generated by RAG Assistant"
    
    for i in range(0, len(lines), chunk_size):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = f"BTEC Content Part {i // chunk_size + 1}"

        body = slide.placeholders[1].text_frame
        body.clear()

        chunk = lines[i:i + chunk_size]

        for line in chunk:
            p = body.add_paragraph()
            p.text = line
            p.level = 0

    file_name = "output.pptx"
    prs.save(file_name)
    return file_name

# =========================================================
# UI
# =========================================================

st.markdown(
    '<div class="title">📚 BTEC RAG Assistant</div>',
    unsafe_allow_html=True
)

question = st.text_area(
    "💬 Ask your question",
    placeholder="Example: What are the assessment criteria for Unit 2 Business?",
    height=120
)

task_type = st.selectbox(
    "Choose content type",
    [
        "Normal Answer",
        "Quiz Questions",
        "Assignment Guidance",
        "Lesson Summary",
        "Learning Activities",
        "Short Explanation"
    ]
)

output_format = st.selectbox(
    "Choose output format",
    ["Text", "PDF", "PowerPoint"]
)

if st.button("Generate content"):
    
    if not question.strip():
        st.warning("Please enter a question")

    else:
        with st.spinner("🔎 Searching approved BTEC sources..."):
            time.sleep(0.3)
            result, ctx = answer(question, task_type)

        if output_format == "Text":
            st.markdown("### Answer")

            formatted_result = result.replace("\n", "<br>")

            st.markdown(
                f'<div class="answer-box">{formatted_result}</div>',
                unsafe_allow_html=True
            )

        elif output_format == "PDF":
            pdf_file = create_pdf(result)

            with open(pdf_file, "rb") as f:
                st.download_button(
                    "Download PDF",
                    f,
                    file_name="btec_output.pdf",
                    mime="application/pdf"
                )

        elif output_format == "PowerPoint":
            ppt_file = create_ppt(result)

            with open(ppt_file, "rb") as f:
                st.download_button(
                    "Download PowerPoint",
                    f,
                    file_name="btec_output.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

        st.markdown("### Sources")

        if ctx:
            source = ctx[0]["source"]
            unit = ctx[0]["unit"]

            st.markdown(f'''
            <div class="source-box">
            📄 {source}<br>
            📌 Unit: {unit}
            </div>
            ''', unsafe_allow_html=True)

   